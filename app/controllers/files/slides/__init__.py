import os
from typing import Any, Dict, List

import pptx

from .schemas import *
from ..exc import BadFile, FileNotFound
from ....db import Session
from ....dependencies import session
from ....settings import STORAGE, STORAGE_PRESENTATION_IMAGES


@session
async def get_slides_info(file_id: int, db: Session) -> List[Slide]:
    filename = await get_slides_filename(file_id, db)
    return get_slides_info_by_filename(file_id, "." + filename.split(".")[-1])


async def get_slides_filename(file_id: int, db: Session) -> str:
    r: Optional[Dict[str, Any]] = await db.fetchrow(
        "SELECT filename FROM files WHERE id=$1 AND filename LIKE '%.ppt%'", file_id
    )

    if r is None:
        raise FileNotFound

    return r["filename"]


def get_slides_info_by_filename(file_id: int, extension: str):
    try:
        p: pptx.Presentation() = pptx.Presentation(
            STORAGE / (str(file_id) + extension),
        )
    except Exception:
        raise BadFile

    return list(
        map(
            lambda x: Slide(
                image_path=os.path.relpath(STORAGE_PRESENTATION_IMAGES, STORAGE)
                + f"/{file_id}-"
                f"{(x[0] + 1 > 9) and (x[0] + 1) or f'0{x[0] + 1}'}.png",
                comment=x[1].has_notes_slide
                and x[1].notes_slide.notes_text_frame.text.strip()
                or "",
            ),
            enumerate(p.slides),
        )
    )
